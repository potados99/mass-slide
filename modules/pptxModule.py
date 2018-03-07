#-*- coding: utf-8 -*-

import re 
import sys 

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from modules import templateModule
from modules import fontModule


def new_presentation():
	return Presentation()

def add_slide(presentation, blank):
	newSlide = prs.slides.add_slide(prs.slide_layouts[6]) 

	left = top = Inches(0)
	width = Inches(10)
	height = Inches(7.5)

	backgroundShape = newSlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

	if blank: COLOR = RGBColor(0, 0, 0)
	else: COLOR = RGBColor(255, 252, 197)

	line = backgroundShape.line
	line.color.rgb = COLOR

	fill = backgroundShape.fill
	fill.solid()
	fill.fore_color.rgb = COLOR

	if not blank:
		continue

	textBox = newSlide.shapes.add_textbox(left, top, width, height)
	text_frame = textBox.text_frame
	text_frame.word_wrap = True
	text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

	return text_frame

def add_line(text_frame, text, size, bold):
	p = text_frame.add_paragraph()
	p.alignment = PP_ALIGN.CENTER

	run = p.add_run()
	run.text = text
	run.font.bold = bold
	run.font.size = Pt(size)
	run.font.name = u"1훈하얀고양이 R"

def write_presentation(Presentation, List):
	for chapter in List:

		linebreakCount = 0
		SIZE = 0
		BOLD = False

		for line in chapter:

			if (chapter.index(line) == 0): # Title
				print ("[" + line + "]")
				if not (chapter[1] == "//"): 
					tf = add_slide(presentation=Presentation, blank=False)

			elif (line == "//"): # Blank mark
				if (linebreakCount >= 1): 
					linebreakCount = 0 
					tf = add_slide(presentation=Presentation, blank=True)

			elif (line == "***"): # Bold mark
				if (BOLD): BOLD = False
				else: BOLD = True

			elif (line == ""): # Linebreak 
				linebreakCount += 1

			else: # Text
				if (linebreakCount >= 1): 
					linebreakCount = 0; 
					tf = add_slide(presentation=Presentation, blank=False)

				boldMatch = re.match('\*\*(.+)\*\*', line)
				tinyMatch = re.match('\'(.+)\'', line)

				SIZE = fontModule.set_font_size(chapter[0])

				if (boldMatch): # Bold
					if (chapter.index(line) == 1): 
						SIZE = fontModule.font_size('big')
					add_line(text_frame=tf, text=boldMatch.group(1), size=SIZE, bold=True)

				elif (tinyMatch): # Tiny
					SIZE = fontModule.font_size('tiny')
					add_line(text_frame=tf, text=tinyMatch.group(1), size=SIZE, bold=False)

				else: # Plane
					add_line(text_frame=tf, text=line, size=SIZE, bold=BOLD)

def save_presentation(prs, path):
	prs.save(path)

'''
prs = Presentation()

List = templateModule.read_template(sys.argv[1])

write_presentation(List)

prs.save('/home/pi/WebDAV/test.pptx')	
'''
