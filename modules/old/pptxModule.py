#-*- coding: utf-8 -*-
from pptx import Presentation 
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import time
import textModule

def makeSlide(slideName, slideLayout, text, bold):
    print('Creating new slide')
    print(slideName)
    newSlide = prs.slides.add_slide(prs.slide_layouts[slideLayout]) 
    #blank layout

    print('Setting size of textbox')
    left = top = Inches(0.1) 
    width = Inches(9.8) 
    height = Inches(7.3)

    print('Creating new textbox')
    textBox = newSlide.shapes.add_textbox(left, top, width, height)

    print('Setting textframe')
    tf = textBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    print('Setting paragraph')
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER

    print('Setting run')
    run = para.add_run()

    print('Setting text and font')
    run.text = text
    run.font.size = Pt(90)
    run.font.bold = bold
    run.font.name = u"1훈하얀고양이 R"

def createNewPresentation():
	prs = Presentation()
	return prs

def addNewSlide(prs, slideLayout):
	newSlide = prs.slides.add_slide(prs.slide_layouts[slideLayout]) 

	return newSlide

def addNewLine():
	return nul

print('Getting list of lines from textModule')
lineList = textModule.makeLineList(textModule.readTextFile())

print('Calling writeToTextFile from textModule')
textModule.writeToTextFile(lineList)

print('Creating presentation')
prs = Presentation()

for i in range(len(lineList)):
    print('Creating new slides')
    makeSlide(slideName='Slide'+str(i+1) , slideLayout=6, text=lineList[i], bold=False)

print('Creating file name')
now = time.localtime() 
currentDateTime = "%04d-%02d-%02d %02d:%02d:%02d" % (now.tm_year, now.tm_mon, now.tm_mday, now.tm_hour, now.tm_min, now.tm_sec)

print('Saving as ' + currentDateTime + '.pptx')
prs.save('../' + currentDateTime + '.pptx')
