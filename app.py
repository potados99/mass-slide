#-*- coding: utf-8 -*-

import re 
import sys 

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from modules import templateModule
from modules import fontModule


def add_slide(blank):
	newSlide = prs.slides.add_slide(prs.slide_layouts[6]) 

	left = top = Inches(0)
	width = Inches(10)
	height = Inches(7.5)

	backgroundShape = newSlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

	if blank: COLOR = RGBColor(0, 0, 0)
	else: COLOR = RGBColor(255, 252, 197)

	line = backgroundShape.line
	line.color.rgb = COLOR

	fill = backgroundShape.fill
	fill.solid()
	fill.fore_color.rgb = COLOR

	if not blank:
		continue

	textBox = newSlide.shapes.add_textbox(left, top, width, height)
	global text_frame
	text_frame = textBox.text_frame
	text_frame.word_wrap = True
	text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
	text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def new_line(text, size, bold):
	p = text_frame.add_paragraph()
	p.alignment = PP_ALIGN.CENTER

	run = p.add_run()
	run.text = text
	run.font.bold = bold
	run.font.size = Pt(size)
	run.font.name = u"1훈하얀고양이 R"

def write_presentation(List):
	for chapter in List:

		linebreakCount = 0
		SIZE = 0
		BOLD = False

		for line in chapter:
			boldMatch = re.match('\*\*.+\*\*', line)
			tinyMatch = re.match('\'.+\'', line)

			if (line == ""): linebreakCount += 1
			# Linebreak
			# 줄바꿈

			elif (line == "***"):
			# Bold mark
			# 볼드체 표시

				if (BOLD): BOLD = False
				else: BOLD = True
				# Toggle
				# 토글

			elif (line == "//"):
			# Empty mark
			# 빈 페이지 표시

				if (linebreakCount >= 1): 
					continue

				linebreakCount = 0 
				add_slide(blank=True)
				# If empty mark seen after linebreak(s),
				# then create new blank slide
					
				# 빈 페이지 표시가 줄바꿈(들) 후에 나타났을 때
				# 새로운 빈 슬라이드를 생성

			elif (chapter.index(line) == 0):
			# Title
			# 제목

				print ("[" + line + "]")

				if (chapter[1] == "//"): 
					continue

				add_slide(blank=False)
				# If the chapter doesn't start with blank slide,
				# then create new slide

				# 챕터가 빈 슬라이드로 시작하지 않을 때
				# 새 슬라이드를 생성
					
			else:
			# Text (bold | tiny | plane)
			# 텍스트 (볼드 | 작음 | 평문)

				SIZE = fontModule.set_font_size(chapter[0])
				# Set font size by chapter title
				# 챕터 제목에 따라 폰트 사이즈 설정

				if (linebreakCount < 1): 
					continue

				linebreakCount = 0; 
				add_slide(blank=False)
				# If (bold | tiny | plane) texts seen after linebreak(s),
				# then add new slide before add new line

				# 텍스트가 줄바꿈(들) 후에 나타났을 때
				# 새 줄을 추가하기 전에 새 슬라이드를 먼저 추가

				if (boldMatch):
				# Bold
				# 볼드체

					if (chapter.index(line) == 1): SIZE = fontModule.font_size('big')
					# If line is the first bold text,
					# then override font size set by chapter title

					# 행이 첫번째 볼드체 텍스트라면,
					# 챕터 제목에 의해 정해진 폰트 사이즈를 변경

					new_line(text=boldMatch.group().replace("*", ""), size=SIZE, bold=True)

				elif (tinyMatch):
				# Tiny
				# 작음

					SIZE = fontModule.font_size('tiny')
					# If line is the 'tiny',
					# override font size set by chapter title

					# 행이 작은 텍스트라면,
					# 챕터 제목에 의해 정해진 폰트 사이즈를 변경

					new_line(text=tinyMatch.group().replace("\'", ""), size=SIZE, bold=False)

				else:
				# Plane
				# 평문

					new_line(text=line, size=SIZE, bold=BOLD)


prs = Presentation()

List = templateModule.delete_space(templateModule.make_chapterList(templateModule.read_template(sys.argv[1])))

write_presentation(List)

prs.save('/home/pi/WebDAV/test.pptx')	
